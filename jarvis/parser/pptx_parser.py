"""
parser/pptx_parser.py
======================

PowerPoint (.pptx) parsing via `python-pptx`.

Design decisions
-----------------
- Each slide is treated as its own "page" (`page_number` = slide index),
  which maps naturally onto the citation model — "slide 7 of
  Failure_Report.pptx" is exactly how an engineer would refer back to it.
- The slide's title placeholder (if any) becomes both a HEADING element and
  the `section_path` for every other element on that slide, so search
  results and citations show which slide's *topic* an answer came from,
  not just its number.
- Speaker notes are extracted as their own BODY_TEXT elements with
  `extra={"is_speaker_note": True}` — the spec doesn't explicitly mention
  speaker notes, but engineering review decks often put the actual
  reasoning/caveats in notes rather than on-slide bullets, so skipping
  them would lose real information. They're tagged distinctly so the
  reasoning engine can weight/cite them differently from on-slide content
  if desired.
- Tables on slides are rendered the same pipe-delimited way as everywhere
  else in the system for consistency.
"""

from __future__ import annotations

from pathlib import Path
from typing import List

from models.document import ExtractionElementType
from parser.base_parser import BaseParser, ParsedDocument, ParserError, RawElement
from utils.logger import get_logger

logger = get_logger(__name__)


class PptxParser(BaseParser):
    @classmethod
    def supported_extensions(cls) -> List[str]:
        return [".pptx"]

    def parse(self, path: Path) -> ParsedDocument:
        from pptx import Presentation

        try:
            presentation = Presentation(str(path))
        except Exception as exc:
            raise ParserError(f"Could not open PPTX {path}: {exc}")

        core_props = presentation.core_properties
        doc = ParsedDocument(
            title=core_props.title or path.stem,
            author=core_props.author,
            page_count=len(presentation.slides._sldIdLst),
        )

        order = 0
        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_title = None

            # Title placeholder first, if present, so section_path is set
            # before we process the rest of the slide's shapes.
            if slide.shapes.title and slide.shapes.title.text.strip():
                slide_title = slide.shapes.title.text.strip()
                doc.elements.append(RawElement(
                    text=slide_title, element_type=ExtractionElementType.HEADING,
                    page_number=slide_index, section_path=slide_title, order_index=order,
                ))
                order += 1

            for shape in slide.shapes:
                try:
                    if shape.has_table:
                        table = shape.table
                        rows_text = [" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows]
                        rendered = "\n".join(rows_text)
                        if rendered.strip():
                            doc.elements.append(RawElement(
                                text=rendered, element_type=ExtractionElementType.TABLE,
                                page_number=slide_index, section_path=slide_title, order_index=order,
                            ))
                            order += 1
                        continue

                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if not text or shape == slide.shapes.title:
                            continue
                        doc.elements.append(RawElement(
                            text=text, element_type=ExtractionElementType.BODY_TEXT,
                            page_number=slide_index, section_path=slide_title, order_index=order,
                        ))
                        order += 1
                except Exception as exc:
                    doc.parser_warnings.append(f"Slide {slide_index}: shape extraction failed ({exc})")
                    logger.warning("PPTX %s slide %d shape failed: %s", path, slide_index, exc)

            # Speaker notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    doc.elements.append(RawElement(
                        text=notes_text, element_type=ExtractionElementType.BODY_TEXT,
                        page_number=slide_index, section_path=slide_title, order_index=order,
                        extra={"is_speaker_note": True},
                    ))
                    order += 1

        return doc
